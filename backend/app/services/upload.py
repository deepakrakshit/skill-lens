from __future__ import annotations

from io import BytesIO

from docx import Document
from fastapi import UploadFile
from pypdf import PdfReader
from pptx import Presentation


def _get_extension(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _parse_pdf(buffer: bytes) -> str:
    reader = PdfReader(BytesIO(buffer))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _parse_docx(buffer: bytes) -> str:
    document = Document(BytesIO(buffer))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _parse_pptx(buffer: bytes) -> str:
    presentation = Presentation(BytesIO(buffer))
    lines: list[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)

    return "\n".join(lines)


def _parse_text(buffer: bytes) -> str:
    return buffer.decode("utf-8", errors="ignore")


async def extract_text_from_upload(file: UploadFile) -> tuple[str, str | None, str | None]:
    file_name = file.filename
    if not file_name:
        return "", None, "No file provided"

    ext = _get_extension(file_name)
    buffer = await file.read()

    if ext == "pdf":
        extracted_text = _parse_pdf(buffer)
    elif ext == "docx":
        extracted_text = _parse_docx(buffer)
    elif ext == "pptx":
        extracted_text = _parse_pptx(buffer)
    elif ext in {"txt", "md"}:
        extracted_text = _parse_text(buffer)
    else:
        return (
            "",
            file_name,
            f"Unsupported file format: .{ext}. Supported: .pdf, .docx, .pptx, .txt, .md",
        )

    cleaned = extracted_text.strip()
    if len(cleaned) < 10:
        return (
            "",
            file_name,
            "Could not extract meaningful text from the file. The file may be corrupted, image-based, or empty.",
        )

    return cleaned, file_name, None